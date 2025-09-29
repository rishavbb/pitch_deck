from pptx import Presentation
from PIL import Image
from typing import Dict, Any, List
import os
import io

class PPTExtractor:
    """Extract text content from PowerPoint pitch decks."""
    
    def __init__(self):
        pass
    
    def extract_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text content and images from a PowerPoint file.
        
        Args:
            file_path (str): Path to the PPT/PPTX file
            
        Returns:
            Dict[str, Any]: Extracted content with metadata, text, and images
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        if not file_path.lower().endswith(('.ppt', '.pptx')):
            raise ValueError("File must be a PowerPoint file (.ppt or .pptx)")
        
        try:
            presentation = Presentation(file_path)
            
            # Extract metadata
            metadata = {
                'file_name': os.path.basename(file_path),
                'file_type': 'PowerPoint',
                'num_slides': len(presentation.slides),
                'file_size': os.path.getsize(file_path)
            }
            
            # Extract text from all slides
            slide_content = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content.append({
                        'slide': slide_num,
                        'content': '\n'.join(slide_text)
                    })
            
            # Combine all text
            full_text = '\n\n'.join([slide['content'] for slide in slide_content])
            
            # Extract images from slides
            images = self._extract_images_from_ppt(presentation)
            
            return {
                'metadata': metadata,
                'slides': slide_content,
                'full_text': full_text,
                'images': images,
                'extraction_success': True
            }
            
        except Exception as e:
            return {
                'metadata': {'file_name': os.path.basename(file_path), 'file_type': 'PowerPoint'},
                'slides': [],
                'full_text': '',
                'images': [],
                'extraction_success': False,
                'error': str(e)
            }
    
    def _extract_images_from_ppt(self, presentation: Presentation) -> List[Image.Image]:
        """
        Extract images from PowerPoint slides.
        
        Args:
            presentation (Presentation): python-pptx Presentation object
            
        Returns:
            List[Image.Image]: List of PIL Image objects
        """
        images = []
        
        try:
            for slide_num, slide in enumerate(presentation.slides, 1):
                for shape in slide.shapes:
                    try:
                        # Check if shape contains an image
                        if hasattr(shape, 'image') and shape.image:
                            # Get image data
                            image_stream = io.BytesIO(shape.image.blob)
                            pil_image = Image.open(image_stream)
                            
                            # Filter out very small images (likely icons/decorations)
                            if pil_image.size[0] > 50 and pil_image.size[1] > 50:
                                images.append(pil_image)
                        
                        # Check for picture shapes
                        elif hasattr(shape, 'shape_type') and shape.shape_type == 13:  # Picture type
                            try:
                                if hasattr(shape, 'image') and shape.image:
                                    image_stream = io.BytesIO(shape.image.blob)
                                    pil_image = Image.open(image_stream)
                                    
                                    if pil_image.size[0] > 50 and pil_image.size[1] > 50:
                                        images.append(pil_image)
                            except Exception:
                                continue
                                
                    except Exception as e:
                        print(f"Warning: Could not extract image from slide {slide_num}: {e}")
                        continue
                        
        except Exception as e:
            print(f"Warning: Could not extract images from PowerPoint: {e}")
        
        return images
